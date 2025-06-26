from flask import Flask, request, jsonify, url_for
from pptx import Presentation
import os
from datetime import datetime
import re

app = Flask(__name__)

@app.route('/generate_certificates', methods=['POST'])
def generate_certificates():
    try:
        data = request.get_json()
        selected_rows = data.get('selectedRows', [])

        if not selected_rows:
            return jsonify({"error": "선택된 행이 없습니다."}), 400

        template_path = os.path.join("static", "certi_template.pptx")
        output_dir = "static"
        os.makedirs(output_dir, exist_ok=True)

        if not os.path.exists(template_path):
            return jsonify({"error": f"템플릿 파일이 존재하지 않습니다: {template_path}"}), 500

        prs = Presentation(template_path)
        today_str = datetime.today().strftime("%y%m%d")

        for row in selected_rows:
            name = row.get("fomat_name", "").strip()
            subject = row.get("upper_subject", "")
            amount = row.get("paid_amount", "")
            period = row.get("period", "")
            date_str = datetime.today().strftime("%Y년 %m월 %d일")

            slide = prs.slides.add_slide(prs.slide_layouts[0])
            try:
                slide.shapes[0].text = name
                slide.shapes[1].text = subject
                slide.shapes[2].text = amount
                slide.shapes[3].text = period
                slide.shapes[4].text = date_str
            except IndexError as e:
                return jsonify({"error": f"텍스트박스 부족: {e}"}), 500

        # 템플릿 슬라이드 제거 (맨 앞 슬라이드)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

        # 첫 번째 row의 이름을 파일명에 사용
        first_name = selected_rows[0].get("fomat_name", "수강증").strip()
        # 파일명 안전하게: 특수문자 제거
        safe_name = re.sub(r'[\\/*?:"<>|]', '', first_name)
        filename = f"수강증_{safe_name}_{today_str}.pptx"

        save_path = os.path.join(output_dir, filename)
        prs.save(save_path)

        file_url = url_for('static', filename=filename, _external=True)
        return jsonify({"message": "수강증 저장 성공", "file_url": file_url})

    except Exception as e:
        return jsonify({"error": str(e)}), 500


if __name__ == '__main__':
    app.run(host="0.0.0.0", port=10000, debug=True)
